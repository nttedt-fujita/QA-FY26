"""slide-structures.mdの座標データからdrawio図を生成するスクリプト"""
import xml.etree.ElementTree as ET
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
import html


def cm(emu_val):
    if emu_val is None:
        return 0
    return round(emu_val / 914400 * 2.54, 1)


def px(cm_val):
    """cmをdrawioのピクセルに変換（1cm ≈ 30px for good visibility）"""
    return round(cm_val * 30)


def get_fill_color(shape):
    """シェイプの塗りつぶし色を取得"""
    try:
        fill = shape.fill
        if fill.type is not None:
            fg = fill.fore_color
            if fg and fg.type is not None:
                rgb = fg.rgb
                if rgb:
                    return f"#{rgb}"
    except:
        pass
    return None


def get_line_color(shape):
    """シェイプの線色を取得"""
    try:
        line = shape.line
        if line.fill and line.fill.type is not None:
            fg = line.fill.fore_color
            if fg and fg.type is not None:
                rgb = fg.rgb
                if rgb:
                    return f"#{rgb}"
    except:
        pass
    return None


def shape_to_style(shape):
    """シェイプからdrawioスタイルを推定"""
    fill = get_fill_color(shape)
    line = get_line_color(shape)

    # デフォルトスタイル
    style_parts = ["rounded=1", "whiteSpace=wrap", "fontSize=10"]

    if fill:
        style_parts.append(f"fillColor={fill}")
    else:
        style_parts.append("fillColor=#FFFFFF")

    if line:
        style_parts.append(f"strokeColor={line}")
    else:
        style_parts.append("strokeColor=#333333")

    return ";".join(style_parts) + ";"


def extract_connectors(slide):
    """スライドからコネクタ（矢印）を抽出"""
    connectors = []
    for shape in slide.shapes:
        # コネクタタイプのシェイプ
        if hasattr(shape, 'begin_x') or shape.shape_type in (
            MSO_SHAPE_TYPE.FREEFORM,
        ):
            continue
        try:
            if shape.shape_type == 9:  # LINE
                connectors.append({
                    "left": cm(shape.left),
                    "top": cm(shape.top),
                    "width": cm(shape.width),
                    "height": cm(shape.height),
                })
        except:
            pass
    return connectors


def slide_to_drawio(pptx_path, slide_num, title):
    """1スライドをdrawio XMLに変換"""
    prs = Presentation(str(pptx_path))
    slide = prs.slides[slide_num - 1]

    # スライドサイズ
    sw = px(cm(prs.slide_width))
    sh = px(cm(prs.slide_height))

    # drawio XML構築
    root = ET.Element("mxGraphModel")
    root.set("dx", "0")
    root.set("dy", "0")
    root.set("grid", "1")
    root.set("gridSize", "10")
    root_cell = ET.SubElement(root, "root")

    # デフォルトセル
    c0 = ET.SubElement(root_cell, "mxCell", id="0")
    c1 = ET.SubElement(root_cell, "mxCell", id="1", parent="0")

    cell_id = 2

    # タイトル
    title_cell = ET.SubElement(root_cell, "mxCell",
        id=str(cell_id),
        value=html.escape(title),
        style="text;fontSize=16;fontStyle=1;align=left;verticalAlign=top;fillColor=none;strokeColor=none;",
        vertex="1",
        parent="1"
    )
    geo = ET.SubElement(title_cell, "mxGeometry",
        x="0", y="0", width=str(sw), height="40",
    )
    geo.set("as", "geometry")
    cell_id += 1

    # Y座標オフセット（タイトル分）
    y_offset = 0

    for shape in slide.shapes:
        # 画像はスキップ
        if shape.shape_type == 13:
            continue

        x = px(cm(shape.left))
        y = px(cm(shape.top)) + y_offset
        w = px(cm(shape.width))
        h = px(cm(shape.height))

        # テキスト取得
        text = ""
        if shape.has_text_frame:
            lines = []
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t:
                    lines.append(html.escape(t))
            text = "<br>".join(lines)

        # グループシェイプ
        if shape.shape_type == 6:
            group_lines = []
            try:
                for child in shape.shapes:
                    if child.has_text_frame:
                        for para in child.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                group_lines.append(html.escape(t))
            except:
                pass
            if group_lines:
                text = "<br>".join(group_lines)

        # テーブル
        if shape.has_table:
            # テーブルは外枠だけ（中身は別途）
            continue

        # テキストもサイズもないシェイプはスキップ
        if not text and w < 10 and h < 10:
            continue

        # スタイル決定
        style = shape_to_style(shape)

        # テキストなしの小さいシェイプ（矢印の線かも）
        if not text and (w < 5 or h < 5):
            # 線として描画
            style = "endArrow=classic;strokeColor=#333333;"

        cell = ET.SubElement(root_cell, "mxCell",
            id=str(cell_id),
            value=text,
            style=style,
            vertex="1",
            parent="1"
        )
        geo = ET.SubElement(cell, "mxGeometry",
            x=str(x), y=str(y), width=str(max(w, 20)), height=str(max(h, 20)),
        )
        geo.set("as", "geometry")
        cell_id += 1

    return ET.tostring(root, encoding="unicode", xml_declaration=False)


# 変換対象（最も重要なもの）
TARGETS = [
    ("群馬通商様_kintone操作マニュアル.pptx", 7, "全体業務フロー（GEN-kintone-群馬通商連携）"),
    ("Kintone経理マニュアル.pptx", 11, "経理業務フロー（kintone導入後）"),
    ("Kintone農業問合せマニュアル.pptx", 3, "農業問合せ業務フロー"),
    ("Kintone受注後マニュアル.pptx", 4, "管理範囲（受注後）"),
    ("Kintone経理マニュアル.pptx", 6, "経理管理範囲"),
    ("Kintone普及管理マニュアル.pptx", 5, "普及部門管理範囲"),
    ("Kintone産業管理マニュアル.pptx", 6, "産業部門管理範囲"),
    ("KintoneSCM業務マニュアル.pptx", 3, "SCM管理範囲"),
]


def main():
    manuals_dir = Path("/home/fuji0130/workspace/QA-FY26/sessions/session316/manuals")
    output_dir = Path("/home/fuji0130/workspace/QA-FY26/sessions/session316/diagrams")
    output_dir.mkdir(exist_ok=True)

    # 個別ファイルとして出力
    for pptx_name, slide_num, title in TARGETS:
        pptx_path = manuals_dir / pptx_name
        safe_name = pptx_path.stem.replace("Kintone", "").replace("kintone", "").replace("様_", "").replace("_", "")
        out_name = f"{safe_name}_s{slide_num:02d}.drawio"
        out_path = output_dir / out_name

        xml_str = slide_to_drawio(pptx_path, slide_num, title)

        # drawioファイルフォーマット
        drawio_xml = f'<?xml version="1.0" encoding="UTF-8"?>\n{xml_str}'
        out_path.write_text(drawio_xml, encoding="utf-8")
        print(f"✅ {out_name} — {title}")

    # 統合ファイル（複数ページ）
    combined_path = output_dir / "kintone-flows-combined.drawio"
    combined_root = ET.Element("mxfile")

    for i, (pptx_name, slide_num, title) in enumerate(TARGETS):
        pptx_path = manuals_dir / pptx_name
        xml_str = slide_to_drawio(pptx_path, slide_num, title)

        diagram = ET.SubElement(combined_root, "diagram",
            id=f"page{i}",
            name=title
        )
        # graphModelをパースしてdiagramの子に追加
        graph_model = ET.fromstring(xml_str)
        diagram.append(graph_model)

    combined_xml = ET.tostring(combined_root, encoding="unicode", xml_declaration=True)
    combined_path.write_text(combined_xml, encoding="utf-8")
    print(f"\n✅ 統合ファイル: {combined_path.name}（{len(TARGETS)}ページ）")


if __name__ == "__main__":
    main()
