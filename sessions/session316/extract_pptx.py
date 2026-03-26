"""kintoneマニュアルPPTXからテキスト・画像情報を抽出するスクリプト"""
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt


def extract_pptx(pptx_path: Path, output_path: Path):
    """PPTXファイルからスライドごとのテキストと画像情報を抽出"""
    prs = Presentation(str(pptx_path))

    lines = []
    lines.append(f"# {pptx_path.stem}")
    lines.append("")
    lines.append(f"- スライド数: {len(prs.slides)}")

    # スライド幅・高さ
    w_cm = prs.slide_width / 914400 * 2.54
    h_cm = prs.slide_height / 914400 * 2.54
    lines.append(f"- スライドサイズ: {w_cm:.1f}cm x {h_cm:.1f}cm")
    lines.append("")
    lines.append("---")
    lines.append("")

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## スライド {i}")
        lines.append("")

        # レイアウト名
        if slide.slide_layout:
            lines.append(f"**レイアウト**: {slide.slide_layout.name}")
            lines.append("")

        # テキスト抽出
        texts = []
        tables_data = []
        image_count = 0

        for shape in slide.shapes:
            # テーブル
            if shape.has_table:
                table = shape.table
                table_rows = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        cell_text = cell.text.strip().replace("\n", " ")
                        row_data.append(cell_text)
                    table_rows.append(row_data)
                tables_data.append(table_rows)
            # テキストフレーム
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # インデントレベルに応じてマーク
                        level = para.level or 0
                        prefix = "  " * level + "- " if level > 0 else ""
                        texts.append(f"{prefix}{text}")
            # 画像
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1

        if texts:
            lines.append("### テキスト")
            lines.append("")
            for t in texts:
                lines.append(t)
            lines.append("")

        if tables_data:
            for ti, table_rows in enumerate(tables_data):
                lines.append(f"### テーブル {ti + 1}")
                lines.append("")
                if table_rows:
                    # ヘッダー
                    header = table_rows[0]
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in table_rows[1:]:
                        # 列数を合わせる
                        while len(row) < len(header):
                            row.append("")
                        lines.append("| " + " | ".join(row[:len(header)]) + " |")
                lines.append("")

        if image_count > 0:
            lines.append(f"**画像**: {image_count}枚")
            lines.append("")

        if not texts and not tables_data and image_count == 0:
            lines.append("（空スライド）")
            lines.append("")

        lines.append("---")
        lines.append("")

    output_path.write_text("\n".join(lines), encoding="utf-8")
    print(f"✅ {pptx_path.name} → {output_path.name} ({len(prs.slides)}スライド)")


def main():
    manuals_dir = Path("/home/fuji0130/workspace/QA-FY26/sessions/session316/manuals")
    output_dir = Path("/home/fuji0130/workspace/QA-FY26/sessions/session316")

    for pptx_file in sorted(manuals_dir.glob("*.pptx")):
        # 出力ファイル名: マニュアル名.md
        safe_name = pptx_file.stem.replace(" ", "_")
        output_path = output_dir / f"{safe_name}.md"
        try:
            extract_pptx(pptx_file, output_path)
        except Exception as e:
            print(f"❌ {pptx_file.name}: {e}")


if __name__ == "__main__":
    main()
